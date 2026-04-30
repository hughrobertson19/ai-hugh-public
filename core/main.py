#!/usr/bin/env python3
import sys
from pathlib import Path
import json
from datetime import datetime, timedelta

import os
from dotenv import load_dotenv
import anthropic

load_dotenv()

client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
sys.path.append(str(Path(__file__).resolve().parent.parent))

# Import path constants from central module
from paths import (
    TASKS_FILE,
    DAILY_METRICS_FILE,
    PERFORMANCE_CONFIG_FILE,
    PERFORMANCE_LOG_FILE,
)

# Podcast generation imports
from typing import Optional
from podcast_generator import PodcastGenerator

# Human-readable due date formatter
def format_due_human(iso_or_text):
    try:
        dt = datetime.fromisoformat(iso_or_text)
        return dt.strftime("%a %b %d @ %I:%M %p")
    except Exception:
        return iso_or_text


def get_due(task):
    """Get the due date from a task, checking both 'due_at' and 'due' fields."""
    return task.get("due_at") or task.get("due")


# `open` = legacy tasks from core/main.py flows. `pending` = tasks synced from
# process_meeting.py. Completed/archived/cancelled are not actionable.
OPEN_STATUSES = {"open", "pending"}


def is_active(task):
    """True if the task counts as actionable today. Missing status defaults to `open`."""
    return isinstance(task, dict) and task.get("status", "open") in OPEN_STATUSES


# -------------------- PODCAST HELPERS --------------------

def extract_text_from_file(file_path: Path) -> str:
    """Extract text from .txt, .pdf, .docx, .pptx. Best-effort v1."""
    suffix = file_path.suffix.lower()

    if suffix in {".txt", ".md"}:
        return file_path.read_text(encoding="utf-8", errors="ignore")

    if suffix == ".pdf":
        try:
            from pypdf import PdfReader  # type: ignore
        except Exception as e:
            raise RuntimeError("Missing dependency for PDF extraction: pypdf. Install with `pip install pypdf`.") from e

        reader = PdfReader(str(file_path))
        parts = []
        for page in reader.pages:
            try:
                parts.append(page.extract_text() or "")
            except Exception:
                parts.append("")
        return "\n".join(parts)

    if suffix == ".docx":
        try:
            from docx import Document  # type: ignore
        except Exception as e:
            raise RuntimeError("Missing dependency for DOCX extraction: python-docx. Install with `pip install python-docx`.") from e

        doc = Document(str(file_path))
        return "\n".join(p.text for p in doc.paragraphs if p.text)

    if suffix == ".pptx":
        try:
            from pptx import Presentation  # type: ignore
        except Exception as e:
            raise RuntimeError("Missing dependency for PPTX extraction: python-pptx. Install with `pip install python-pptx`.") from e

        prs = Presentation(str(file_path))
        slides_text = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_parts = [f"[Slide {i}]" ]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    t = shape.text.strip()
                    if t:
                        slide_parts.append(t)
            slides_text.append("\n".join(slide_parts))
        return "\n\n".join(slides_text)

    raise ValueError(f"Unsupported file type: {suffix}. Use .txt, .pdf, .docx, or .pptx")



def parse_flag_value(flag: str) -> Optional[str]:
    """Return the value after a flag like --title, or None."""
    if flag not in sys.argv:
        return None
    idx = sys.argv.index(flag)
    if idx + 1 >= len(sys.argv):
        return None
    return sys.argv[idx + 1]


# --- Podcast auto-length estimator ---
def estimate_minutes_from_text(text: str) -> int:
    """Estimate a sensible episode length from source text.

    Heuristic: ~165 spoken words/minute, with guardrails.
    Returns minutes in [3, 15].
    """
    words = len((text or "").split())
    if words <= 0:
        return 5
    est = int(round(words / 165))
    return max(3, min(15, est))



def classify_task_priority(task):
    """
    Minimal v1 priority classifier for Daily Brief and Today views.
    """
    due_raw = get_due(task)
    if not due_raw:
        return "IGNORE"

    try:
        due_dt = datetime.fromisoformat(due_raw)
    except Exception:
        # Fallback for legacy text-based due values
        due_text = str(due_raw).lower()
        if "now" in due_text or "asap" in due_text:
            return "DO NOW"
        if "today" in due_text:
            return "DO TODAY"
        return "IGNORE"

    now = datetime.now()

    # Overdue or within next 2 hours
    if due_dt <= now + timedelta(hours=2):
        return "DO NOW"

    # Due later today
    if due_dt.date() == now.date():
        return "DO TODAY"

    return "IGNORE"


# Helper: determine if a task is deal-related
def is_deal_related(task):
    """
    Determine whether a task is directly associated with closing a deal.
    v1 heuristic: Calls and Messages are considered deal-related.
    """
    return task.get("action_type") in {"Call", "Message"}


def previous_business_day_at_9am(dt):
    """
    Minimal stub: return previous day at same time.
    Improved business-day logic can be added later if needed.
    """
    return dt - timedelta(days=1)


# Helper: add business days at 9:00 AM
def add_business_days_at_9am(start_dt, business_days):
    """
    Return a datetime that is `business_days` business days after start_dt,
    set to 9:00 AM local time.
    Weekends are skipped.
    """
    current = start_dt
    added = 0

    while added < business_days:
        current += timedelta(days=1)
        if current.weekday() < 5:  # Mon–Fri
            added += 1

    return current.replace(hour=9, minute=0, second=0, microsecond=0)


def business_days_since(dt):
    """
    Return number of business days between dt and now.
    Weekends are excluded.
    """
    now = datetime.now()
    if dt > now:
        return 0

    days = 0
    current = dt

    while current.date() < now.date():
        current += timedelta(days=1)
        if current.weekday() < 5:  # 0=Mon, 6=Sun
            days += 1

    return days

def load_tasks():
    if not TASKS_FILE.exists():
        return {}
    data = json.loads(TASKS_FILE.read_text())
    if isinstance(data, list):
        # migrate old list format to dict
        migrated = {}
        for idx, task in enumerate(data):
            migrated[str(idx)] = task
        return migrated
    return data

def save_task(lead_id, task):
    """
    Enforce single active task per lead.
    """
    tasks = load_tasks()
    tasks[str(lead_id)] = task
    TASKS_FILE.write_text(json.dumps(tasks, indent=2))


# --- DAILY METRICS HELPERS ---
def load_daily_metrics():
    if not DAILY_METRICS_FILE.exists():
        return {}
    return json.loads(DAILY_METRICS_FILE.read_text())


def increment_metric(metric_name):
    today = datetime.now().strftime("%Y-%m-%d")
    data = load_daily_metrics()

    if today not in data:
        data[today] = {
            "calls": 0,
            "messages": 0,
            "new_prospects": 0,
            "meetings_advanced": 0
        }

    data[today][metric_name] += 1
    DAILY_METRICS_FILE.write_text(json.dumps(data, indent=2))


# Coaching log function
def save_coaching_log(entry):
    """
    Append coaching feedback to coaching_log.json.
    Keeps performance analytics separate from execution tasks.
    """
    coaching_file = Path(__file__).parent / "coaching_log.json"

    if coaching_file.exists():
        try:
            existing = json.loads(coaching_file.read_text())
            if not isinstance(existing, list):
                existing = []
        except Exception:
            existing = []
    else:
        existing = []

    existing.append(entry)
    coaching_file.write_text(json.dumps(existing, indent=2))

# --- PERFORMANCE / SCOREBOARD LOGIC ---
def load_performance_config():
    if not PERFORMANCE_CONFIG_FILE.exists():
        return {
            "weights": {
                # Volume Layer (40%)
                "calls": 1,
                "messages": 1,
                "new_prospects": 2,

                # Leverage Layer (60%)
                "meetings_advanced": 5,
                "multi_thread_bonus": 4,
                "high_intent_bonus": 3,
                "calendar_lock_bonus": 3,
                "anchored_next_step_bonus": 2
            },
            "daily_target_score": 15
        }
    return json.loads(PERFORMANCE_CONFIG_FILE.read_text())


def log_daily_performance(score):
    today = datetime.now().strftime("%Y-%m-%d")

    if PERFORMANCE_LOG_FILE.exists():
        try:
            data = json.loads(PERFORMANCE_LOG_FILE.read_text())
            if not isinstance(data, dict):
                data = {}
        except Exception:
            data = {}
    else:
        data = {}

    data[today] = {
        "score": score
    }

    PERFORMANCE_LOG_FILE.parent.mkdir(parents=True, exist_ok=True)
    PERFORMANCE_LOG_FILE.write_text(json.dumps(data, indent=2))


def calculate_today_score():
    today = datetime.now().strftime("%Y-%m-%d")
    metrics = load_daily_metrics()
    config = load_performance_config()
    tasks = load_tasks()

    weights = config.get("weights", {})

    score = 0

    # -------------------------
    # VOLUME LAYER
    # -------------------------
    if today in metrics:
        today_metrics = metrics[today]
        for key in ["calls", "messages", "new_prospects"]:
            weight = weights.get(key, 0)
            score += today_metrics.get(key, 0) * weight

    # -------------------------
    # LEVERAGE LAYER
    # -------------------------
    for _, task in tasks.items():

        # Meetings advanced
        if task.get("stage") == "Post-Meeting":
            score += weights.get("meetings_advanced", 0)

        # Multi-thread bonus (2+ stakeholders detected)
        stakeholders = task.get("stakeholders_detected", [])
        if isinstance(stakeholders, list) and len(stakeholders) >= 2:
            score += weights.get("multi_thread_bonus", 0)

        # High intent bonus
        if task.get("intent_score", 0) >= 70:
            score += weights.get("high_intent_bonus", 0)

        # Calendar lock bonus (meeting scheduled with datetime)
        if get_due(task) and "T" in get_due(task):
            score += weights.get("calendar_lock_bonus", 0)

        # Anchored next step bonus (explicit next step title)
        if task.get("title") and "follow" not in task.get("title", "").lower():
            score += weights.get("anchored_next_step_bonus", 0)

    return score

def record_call_attempt(lead_id):
    tasks = load_tasks()

    if str(lead_id) not in tasks:
        print("No active task for this lead.")
        return

    task = tasks[str(lead_id)]

    task["attempts"] = task.get("attempts", 0) + 1
    task["last_attempt_at"] = datetime.now().isoformat()

    task["action_type"] = "Message"
    task["due_at"] = datetime.now().replace(hour=10, minute=0, second=0, microsecond=0).isoformat()
    task["title"] = "Follow up after missed call"

    tasks[str(lead_id)] = task
    TASKS_FILE.write_text(json.dumps(tasks, indent=2))

    increment_metric("calls")

    print(
        f"Logged call attempt #{task['attempts']} — "
        f"task updated to MESSAGE follow-up for lead '{lead_id}'."
    )

def refine_task_for_context(output, task):
    """
    Adjust task title, action type, and due timing based on momentum
    and meeting context (enterprise-style refinement).
    """

    if output.momentum_status and "High" in output.momentum_status:
        task["action_type"] = "Call"
        task["title"] = "Align on agenda and confirm meeting details"

        # If we have a real meeting datetime, set due = previous business day @ 9:00 AM
        if getattr(output, "meeting_datetime", None):
            due_dt = previous_business_day_at_9am(output.meeting_datetime)
            task["due_at"] = due_dt.isoformat()
        # else: keep whatever due was set upstream as a fallback

    return task

def main():
    # --- Unified morning view (delegates to today.py at project root) ---
    if "--morning" in sys.argv:
        import runpy
        today_path = Path(__file__).resolve().parent.parent / "today.py"
        if not today_path.exists():
            print(f"today.py not found at {today_path}")
            return
        # Run today.py as __main__ so its `if __name__ == "__main__":` block fires.
        runpy.run_path(str(today_path), run_name="__main__")
        return

    # --- Podcast generation (v1) ---
    if "--make-podcast" in sys.argv:
        idx = sys.argv.index("--make-podcast")
        if idx + 1 >= len(sys.argv):
            print("Usage: python3 core/main.py --make-podcast <file.pdf|file.pptx|file.docx|file.txt> [--title <title>] [--minutes <n>] [--no-tts]")
            return

        input_path = Path(sys.argv[idx + 1]).expanduser()
        if not input_path.exists():
            raise FileNotFoundError(f"Input file not found: {input_path}")

        title = parse_flag_value("--title") or input_path.stem
        minutes_raw = parse_flag_value("--minutes")
        tts = ("--no-tts" not in sys.argv)

        source_text = extract_text_from_file(input_path)

        if minutes_raw:
            minutes = int(minutes_raw)
            minutes_note = "(forced)"
        else:
            minutes = estimate_minutes_from_text(source_text)
            minutes_note = "(auto)"

        print(f"Generating podcast: '{title}' (~{minutes} min) {minutes_note} from {input_path.name} ...")

        gen = PodcastGenerator(enable_tts=tts, default_minutes=minutes)
        result = gen.make_episode(title=title, source_text=source_text, minutes=minutes, tts=tts)

        print("\nDONE")
        print("Episode dir:", result.episode_dir)
        print("Script:", result.script_path)
        print("Manifest:", result.manifest_path)
        if result.audio_path:
            print("Audio:", result.audio_path)
        else:
            print("Audio: (skipped)")
        return
    # --- Daily Brief (read-only) ---
    if "--brief" in sys.argv:
        print("GOOD MORNING, HUGH\n")

        tasks = load_tasks()

        # Prioritize deal-related first
        prioritized = {
            "DO NOW": [],
            "DO TODAY": [],
            "IGNORE": []
        }

        for lead, task in tasks.items():
            if not is_active(task):
                continue
            bucket = classify_task_priority(task)
            if is_deal_related(task):
                prioritized[bucket].append((lead, task))

        selected = None
        if prioritized["DO NOW"]:
            selected = prioritized["DO NOW"][0]
        elif prioritized["DO TODAY"]:
            selected = prioritized["DO TODAY"][0]

        # If no deal-related task exists, fall back to non-deal tasks
        if not selected:
            non_deal_prioritized = {
                "DO NOW": [],
                "DO TODAY": [],
                "IGNORE": []
            }

            for lead, task in tasks.items():
                if not is_active(task):
                    continue
                if is_deal_related(task):
                    continue
                bucket = classify_task_priority(task)
                non_deal_prioritized[bucket].append((lead, task))

            if non_deal_prioritized["DO NOW"]:
                selected = non_deal_prioritized["DO NOW"][0]
            elif non_deal_prioritized["DO TODAY"]:
                selected = non_deal_prioritized["DO TODAY"][0]

        # DAILY BRIEF — ENFORCED FORMAT
        if selected:
            lead, task = selected
            print("TODAY’S NON-NEGOTIABLE")
            print(f"[{task['action_type']}] {task['title']}")
            print(f"Due: {format_due_human(get_due(task))}")

            urgency = task.get("urgency_score", 0)
            intent = task.get("intent_score", 0)
            stakeholders = task.get("stakeholders_detected", [])

            print(f"Urgency Score: {urgency}")
            print(f"Intent Score: {intent}")
            print(f"Stakeholders Detected: {len(stakeholders)}\n")

            print("WHY THIS COMES FIRST")
            if is_deal_related(task):
                print("Protects revenue and deal momentum.\n")
            else:
                print("This task is time-sensitive and creates downstream risk if delayed.\n")

            # Push line logic
            if urgency >= 90:
                print("EXECUTION PUSH: Do this immediately. Protect the deal.\n")
            elif urgency >= 75:
                print("EXECUTION PUSH: Lock this in before 10am.\n")
            else:
                print("EXECUTION PUSH: Complete before moving to pipeline building.\n")

            print("AFTER THIS")
            print("Prospect for new business.\n")
        else:
            print("TODAY’S NON-NEGOTIABLE")
            print("None.\n")

            print("WHY THIS COMES FIRST")
            print("No deadlines or deal risk detected today.\n")

            print("AFTER THIS")
            print("Prospect for new business.\n")

        return

    # --- Today view (smart execution brain) ---
    if "--today" in sys.argv:
        tasks = load_tasks()
        now = datetime.now()

        prioritized = {
            "DO NOW": [],
            "DO TODAY": [],
            "IGNORE": []
        }

        for lead, task in tasks.items():
            if not is_active(task):
                continue
            bucket = classify_task_priority(task)
            prioritized[bucket].append((lead, task))

        # Sort each time bucket by urgency_score (highest first)
        for bucket in ["DO NOW", "DO TODAY"]:
            prioritized[bucket].sort(
                key=lambda x: x[1].get("urgency_score", 0),
                reverse=True
            )

        print("DAILY EXECUTION PLAN — HUGH\n")

        # Collect HANDLE TODAY tasks
        handle_today = []

        # 1. DO NOW + DO TODAY
        for bucket in ["DO NOW", "DO TODAY"]:
            handle_today.extend(prioritized[bucket])

        # 2. Inactivity trigger (≥ 3 business days, Active stages)
        inactivity_triggers = []
        for lead, task in prioritized["IGNORE"]:
            stage = task.get("stage")
            last_activity_raw = task.get("last_activity_at")

            if stage in {"Active Outreach", "Post-Meeting"} and last_activity_raw:
                try:
                    last_activity_dt = datetime.fromisoformat(last_activity_raw)
                    if business_days_since(last_activity_dt) >= 3:
                        inactivity_triggers.append((lead, task))
                except Exception:
                    pass

        # Sort inactivity triggers by urgency_score (highest first)
        inactivity_triggers.sort(
            key=lambda x: x[1].get("urgency_score", 0),
            reverse=True
        )

        handle_today.extend(inactivity_triggers)

        # 3. High-urgency future tasks (urgency ≥ 90)
        high_risk_future = [
            (lead, task)
            for lead, task in prioritized["IGNORE"]
            if task.get("urgency_score", 0) >= 90
            and (lead, task) not in inactivity_triggers
        ]

        handle_today.extend(high_risk_future)

        # Final sort for DO TODAY logic:
        # First item (if any) will become DO NOW.
        # Remaining items sorted by earliest due date, then urgency_score (desc).
        def sort_key(item):
            lead, task = item

            due_raw = get_due(task)
            try:
                due_dt = datetime.fromisoformat(due_raw)
            except Exception:
                due_dt = datetime.max

            urgency = task.get("urgency_score", 0)

            # Determine priority tier
            # 0 = Due Today
            # 1 = Inactivity Trigger
            # 2 = Due Tomorrow
            # 3 = Everything else
            tier = 3

            now = datetime.now()

            # Overdue escalation (hard override)
            if due_dt < now:
                return (0, due_dt, -100)

            # Due Today
            if due_dt.date() == now.date():
                tier = 0

            # Due Tomorrow
            elif due_dt.date() == (datetime.now().date() + timedelta(days=1)):
                tier = 2

            # Inactivity Trigger (Active stages only, ≥ 3 business days)
            stage = task.get("stage")
            last_activity_raw = task.get("last_activity_at")
            if stage in {"Active Outreach", "Post-Meeting"} and last_activity_raw:
                try:
                    last_activity_dt = datetime.fromisoformat(last_activity_raw)
                    if business_days_since(last_activity_dt) >= 3:
                        tier = 1
                except Exception:
                    pass

            return (tier, due_dt, -urgency)

        handle_today = sorted(handle_today, key=sort_key)

        # Flattened ordered list already stored in handle_today
        print("DO NOW")
        if handle_today:
            lead, task = handle_today[0]
            print(f"- [{task['action_type']}] {task['title']} | Lead: {lead}")
            print(f"  Due: {format_due_human(get_due(task))}")
        else:
            print("- Prospect for new business.")

        print()
        print("DO TODAY (TOP 3 REVENUE IMPACT TASKS)")

        if len(handle_today) > 1:
            top_today = handle_today[1:4]  # Limit to top 3 after DO NOW
            backlog = handle_today[4:]

            for lead, task in top_today:
                print(f"- [{task['action_type']}] {task['title']} | Lead: {lead}")
                print(f"  Due: {format_due_human(get_due(task))}")

            if backlog:
                print()
                print("BACKLOG (DO NOT TOUCH UNTIL TOP 3 COMPLETE)")
                for lead, task in backlog:
                    print(f"- [{task['action_type']}] {task['title']} | Lead: {lead}")
                    print(f"  Due: {format_due_human(get_due(task))}")
        else:
            print("- None")

        print()
        print("2. BUILD PIPELINE")
        print("- Prospect for new business.")

        print()
        print("3. IF TIME REMAINS")
        print("- Work on AI Hugh / career advancement tasks.")

        print()

        return

    # --- End of Day Reality Check ---
    if "--eod" in sys.argv:
        tasks = load_tasks()

        tasks = {lead: t for lead, t in tasks.items() if is_active(t)}
        outstanding = len(tasks)

        # Simple proxy signals (v1)
        pipeline_moved = False
        prospecting_done = False

        for _, task in tasks.items():
            if task.get("action_type") in {"Call", "Message"}:
                prospecting_done = True
            if "meeting" in task.get("title", "").lower():
                pipeline_moved = True

        print("END OF DAY CHECK — HUGH\n")

        print("PIPELINE")
        print(f"• Prospecting done today: {'YES' if prospecting_done else 'NO'}")
        print(f"• New pipeline created: {'YES' if pipeline_moved else 'NO'}\n")

        print("TASKS")
        print(f"• Tasks outstanding: {outstanding}")

        if outstanding > 0:
            print("• Outstanding task(s):")
            for lead, task in tasks.items():
                title = task.get("title", "Untitled task")
                due = get_due(task) or "No due date"
                due_lower = str(due).lower()

                line = f"  - [{task.get('action_type')}] {title} (Due: {due})"

                if "today" in due_lower or "now" in due_lower or "asap" in due_lower:
                    line += "  ← DO THIS BEFORE YOU LEAVE WORK"

                print(line)
        print()

        print("REFLECTION")
        if pipeline_moved or prospecting_done:
            print("Today moved the needle. Keep compounding.\n")
        else:
            print("Today did not materially move pipeline.\n")

        print("Tomorrow’s default priority will be prospecting.\n")

        # --- PERFORMANCE SCORE OUTPUT ---
        score = calculate_today_score()
        config = load_performance_config()
        target = config.get("daily_target_score", 10)

        print("PERFORMANCE SCORE")
        print(f"• Today’s weighted score: {score}")
        print(f"• Target score: {target}")

        # --- WEEKDAY LOCK LOGIC ---
        today_dt = datetime.now()
        is_weekday = today_dt.weekday() < 5  # Mon–Fri only

        metrics = load_daily_metrics()
        today_str = today_dt.strftime("%Y-%m-%d")
        today_metrics = metrics.get(today_str, {})

        meaningful_action = (
            today_metrics.get("new_prospects", 0) > 0 or
            today_metrics.get("meetings_advanced", 0) > 0
        )

        if score >= target:
            print("• Status: ELITE EXECUTION\n")

        elif is_weekday and not meaningful_action:
            print("\nEXECUTION INCOMPLETE.")
            print("Log one meaningful revenue action before shutdown.")

            override = input("Override? (emergency / sick / justified / no) \n> ").strip().lower()

            if override in {"emergency", "sick", "justified"}:
                override_file = PERFORMANCE_LOG_FILE.parent / "overrides.json"
                if override_file.exists():
                    try:
                        override_data = json.loads(override_file.read_text())
                        if not isinstance(override_data, list):
                            override_data = []
                    except Exception:
                        override_data = []
                else:
                    override_data = []

                override_data.append({
                    "date": today_str,
                    "reason": override,
                    "score": score
                })

                override_file.write_text(json.dumps(override_data, indent=2))
                print("Override logged. Shutdown permitted.\n")
            else:
                print("Action required before shutdown.\n")


        else:
            print("• Status: Standard execution — improve leverage tomorrow.\n")

        return

    if "--lead" not in sys.argv:
        # --next-move is the only flag that works without --lead
        if "--next-move" not in sys.argv:
            print("Usage: python3 main.py --lead <lead_id> <notes.txt>")
            return
        lead_idx = -1
        lead_id = "unknown"
    else:
        lead_idx = sys.argv.index("--lead")
        try:
            lead_id = sys.argv[lead_idx + 1]
        except IndexError:
            print("Usage: python3 main.py --lead <lead_id> <notes.txt>")
            return

    if "--attempt" in sys.argv:
        record_call_attempt(lead_id)
        return

    if "--new-prospect" in sys.argv:
        now = datetime.now()

        task = {
            "action_type": "Call",
            "title": "Initial outreach to new prospect",
            "due_at": add_business_days_at_9am(now, 0).isoformat(),
            "stage": "Active Outreach",
            "created_at": now.isoformat(),
            "last_activity_at": now.isoformat(),
            "urgency_score": 75,
            "attempts": 0,
            "intent_score": 0,
            "stakeholders_detected": []
        }

        save_task(lead_id, task)

        increment_metric("new_prospects")

        print(f"New prospect '{lead_id}' created.")
        print("Initial outreach task added for today.\n")
        return

    if "--post-meeting" in sys.argv:
        tasks = load_tasks()
        now = datetime.now()

        follow_up_due = add_business_days_at_9am(now, 2)

        task = {
            "action_type": "Message",
            "title": "Post-meeting follow-up and confirm next step",
            "due_at": follow_up_due.isoformat(),
            "stage": "Post-Meeting",
            "last_activity_at": now.isoformat(),
            "urgency_score": 85,
            "attempts": 0
        }

        save_task(lead_id, task)

        increment_metric("meetings_advanced")

        print(f"Post-meeting follow-up task created for lead '{lead_id}'.")
        print(f"Due: {format_due_human(task['due_at'])}")
        return

    if "--discovery" in sys.argv:
        now = datetime.now()

        print("\nDISCOVERY CALL CAPTURE — HUGH\n")

        next_step = input("What was the agreed next step? \n> ")
        meeting_scheduled = input("Was a follow-up meeting scheduled? (y/n) \n> ").strip().lower()

        if meeting_scheduled == "y":
            meeting_date_str = input("Enter meeting date/time (YYYY-MM-DD HH:MM) \n> ")
            try:
                meeting_dt = datetime.strptime(meeting_date_str, "%Y-%m-%d %H:%M")
                due_dt = previous_business_day_at_9am(meeting_dt)
            except Exception:
                print("Invalid date format. Defaulting to 2 business days.")
                due_dt = add_business_days_at_9am(now, 2)
        else:
            due_dt = add_business_days_at_9am(now, 2)

        task = {
            "action_type": "Message",
            "title": next_step if next_step else "Follow up from discovery call",
            "due_at": due_dt.isoformat(),
            "stage": "Active Outreach",
            "last_activity_at": now.isoformat(),
            "urgency_score": 80,
            "attempts": 0
        }

        save_task(lead_id, task)

        print(f"\nDiscovery follow-up task created for lead '{lead_id}'.")
        print(f"Due: {format_due_human(task['due_at'])}\n")
        return

    if "--next-move" in sys.argv:
        """
        Usage:
            python3 core/main.py <lead_id> --next-move [path/to/note.txt]
            python3 core/main.py <lead_id> --next-move   (paste mode)

        Accepts any sales interaction — email, call notes, meeting note, deal update.
        Runs note_processor for structured extraction, then asks Claude for
        the single best next move and generates the exact message to send.

        Does NOT require the lead to be in tasks.json already.
        """
        sys.path.insert(0, str(Path(__file__).resolve().parent))
        from note_processor import process as tp_process

        # --- Read input ---
        raw_text = ""
        next_move_idx = sys.argv.index("--next-move")
        if len(sys.argv) > next_move_idx + 1 and not sys.argv[next_move_idx + 1].startswith("--"):
            note_path = Path(sys.argv[next_move_idx + 1]).expanduser()
            if not note_path.exists():
                print(f"File not found: {note_path}", file=sys.stderr)
                sys.exit(1)
            raw_text = note_path.read_text(encoding="utf-8", errors="replace")
        else:
            print("Paste the interaction (email, call notes, or meeting note).")
            print("Press Enter twice to finish:\n")
            lines = []
            blank_count = 0
            while True:
                line = input()
                if line == "":
                    blank_count += 1
                    if blank_count >= 2:
                        break
                    continue
                blank_count = 0
                lines.append(line)
            raw_text = "\n".join(lines)

        if not raw_text.strip():
            print("No input provided.", file=sys.stderr)
            sys.exit(1)

        # --- Stage 1: deterministic extraction via note_processor ---
        extracted = tp_process(raw_text)
        meeting_type = extracted["meeting_info"].get("meeting_type", "external")
        summary_block = extracted.get("meeting_summary", {})
        hugh_owns = extracted["commitments"].get("hugh_owns", [])
        others_owe = extracted["commitments"].get("others_owe_hugh", [])
        revenue_signals = extracted.get("revenue_signals", [])
        gaps = extracted.get("follow_up_gaps", [])
        confidence = extracted.get("extraction_confidence", 0.0)

        # --- Stage 2: Claude for next move + message ---
        situation_summary = f"""
Meeting type: {meeting_type}
What happened: {summary_block.get("what_happened", "Unknown")}
Why it matters: {summary_block.get("why_it_matters", "Unknown")}
Extraction confidence: {confidence}

Hugh's commitments:
{chr(10).join(f"  - {t['task']} (due: {t.get('due') or 'no deadline'})" for t in hugh_owns) or "  None extracted"}

Others owe Hugh:
{chr(10).join(f"  - {t['owner']}: {t['task']}" for t in others_owe) or "  None extracted"}

Revenue signals:
{chr(10).join(f"  - Signal: {s.get('signal')} | Risk: {s.get('risk')}" for s in revenue_signals) or "  None extracted"}

Follow-up gaps (unclear/missing):
{chr(10).join(f"  - [{g['issue']}] {g['context']}" for g in gaps[:5]) or "  None"}
"""

        prompt = f"""You are a senior enterprise SDR coach working with Hugh Robertson, an SDR selling enterprise sustainability and product compliance software.

Hugh just had this sales interaction. Here is the structured extraction:

{situation_summary}

Raw interaction text:
---
{raw_text[:3000]}
---

Your job:
1. Write a 2-line SUMMARY of what just happened and where the deal stands.
2. List up to 3 RISKS — things that could cause this to stall or go cold.
3. List up to 3 MISSING pieces — decision process, urgency, stakeholders, problem clarity.
4. State the SINGLE BEST NEXT MOVE (action + channel + timing).
5. Write the EXACT MESSAGE to send — short, direct, ready to copy-paste. No subject line needed unless it's an email. No fluff.

Rules:
- Be specific to this interaction. No generic SDR advice.
- The message must include a clear next step.
- Max 80 words for the message.
- If no clear deal signal exists, say so directly.

Return STRICT JSON:
{{
  "summary": "2-line summary",
  "risks": ["risk 1", "risk 2", "risk 3"],
  "missing": ["gap 1", "gap 2", "gap 3"],
  "next_move": {{
    "action": "what to do",
    "channel": "email | text | call | linkedin",
    "timing": "now | today | tomorrow | this week"
  }},
  "message": "ready-to-send message text"
}}"""

        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=1024,
            system="You are a senior enterprise SDR coach. Return only valid JSON.",
            messages=[
                {"role": "user", "content": prompt}
            ],
            temperature=0.2
        )

        content = response.content[0].text.strip()

        # Robust JSON parse
        parsed = None
        try:
            parsed = json.loads(content)
        except Exception:
            pass
        if parsed is None:
            try:
                start = content.find("{")
                end = content.rfind("}")
                if start != -1 and end != -1:
                    parsed = json.loads(content[start:end + 1])
            except Exception:
                pass

        if parsed is None:
            print("\nCould not parse Claude response. Raw output:\n")
            print(content)
            sys.exit(1)

        # --- Output ---
        SEP = "─" * 56
        print(f"\n{SEP}")
        print("  AI HUGH — NEXT MOVE ENGINE")
        print(SEP)

        print("\nSUMMARY:")
        print(f"  {parsed.get('summary', 'No summary.')}")

        risks = parsed.get("risks", [])
        if risks:
            print("\nRISKS:")
            for r in risks:
                print(f"  ✗  {r}")

        missing = parsed.get("missing", [])
        if missing:
            print("\nMISSING:")
            for m in missing:
                print(f"  ?  {m}")

        nm = parsed.get("next_move", {})
        print("\nNEXT MOVE:")
        print(f"  {nm.get('action', '—')}  |  {nm.get('channel', '—')}  |  {nm.get('timing', '—')}")

        print("\nMESSAGE:")
        print("  " + "\n  ".join(parsed.get("message", "No message generated.").split("\n")))

        print(f"\n{SEP}\n")
        return

    if "--auto-discovery" in sys.argv:
        now = datetime.now()

        # Load meeting notes from file arg if provided, else paste mode
        notes = ""
        if len(sys.argv) > lead_idx + 2 and not sys.argv[lead_idx + 2].startswith("--"):
            notes_path = Path(sys.argv[lead_idx + 2]).expanduser()
            if not notes_path.exists():
                raise FileNotFoundError(f"Notes file not found: {notes_path}")
            notes = notes_path.read_text()
        else:
            print("Paste discovery notes (press Enter twice consecutively to finish):")
            lines = []
            blank_count = 0
            while True:
                line = input()
                if line == "":
                    blank_count += 1
                    if blank_count >= 2:
                        break
                    continue
                blank_count = 0
                lines.append(line)
            notes = "\n".join(lines)

        # Use LLM to extract structured next step
        prompt = f"""
You are an enterprise SDR strategist.

Your job is NOT just to extract what was said.
Your job is to determine the optimal next commercial action to maintain or accelerate deal momentum.

Rules:
- If a clear next step was explicitly agreed, use it.
- If the next step was vague (e.g. "we'll review internally"), propose a concrete follow-up action.
- If the prospect says they need to "review internally" and NO follow-up date is locked, you MUST recommend setting a calendar hold within 3–5 business days.
- If no follow-up meeting is scheduled, treat this as a closing weakness.
- Do NOT invent meetings or commitments that contradict the notes.
- If a meeting datetime is explicitly mentioned, extract it in ISO format.
- Momentum =
    High: Clear urgency, stakeholder alignment, meeting scheduled.
    Medium: Interest shown but internal review required.
    Low: Hesitation, delays, or lack of clarity.
- Intent score (0–100) should increase when:
    • Multiple stakeholders are mentioned
    • Economic buyer or procurement is referenced
    • Technical validation is scheduled
    • Clear timeline urgency exists
  Lower score if:
    • Only one stakeholder
    • Vague next step
    • No timeline anchor

Return STRICT JSON only with:
{{
  "next_step": "...",
  "meeting_datetime": "ISO or null",
  "momentum": "High | Medium | Low",
  "close_strength": "Strong | Moderate | Weak",
  "calendar_hold_recommended": true | false,
  "intent_score": 0,
  "stakeholder_roles_detected": [],
  "close_feedback": "1-2 sentence coaching insight"
}}

Notes:
{notes}
"""

        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=1024,
            system="You extract structured sales follow-up data.",
            messages=[
                {"role": "user", "content": prompt}
            ],
            temperature=0.0
        )

        content = response.content[0].text

        # --- Robust JSON extraction ---
        parsed = None

        # Try direct JSON parse first
        try:
            parsed = json.loads(content)
        except Exception:
            pass

        # If direct parse fails, try to extract first JSON object from text
        if parsed is None:
            try:
                start = content.find("{")
                end = content.rfind("}")
                if start != -1 and end != -1:
                    json_block = content[start:end+1]
                    parsed = json.loads(json_block)
            except Exception:
                pass

        # Final fallback
        if parsed is None:
            parsed = {
                "next_step": "Follow up from discovery call",
                "meeting_datetime": None,
                "momentum": "Medium"
            }

        next_step = parsed.get("next_step", "Follow up from discovery call")
        close_strength = parsed.get("close_strength", "Unknown")
        close_feedback = parsed.get("close_feedback", "No coaching feedback available.")
        calendar_hold_recommended = parsed.get("calendar_hold_recommended", False)
        intent_score = parsed.get("intent_score", 50)
        stakeholder_roles = parsed.get("stakeholder_roles_detected", [])
        meeting_dt = None

        if parsed.get("meeting_datetime"):
            try:
                meeting_dt = datetime.fromisoformat(parsed["meeting_datetime"])
            except Exception:
                meeting_dt = None

        due_dt = add_business_days_at_9am(now, 2)
        if meeting_dt:
            due_dt = previous_business_day_at_9am(meeting_dt)

        urgency_score = 80
        if parsed.get("momentum") == "Low":
            urgency_score = 90
        elif parsed.get("momentum") == "High":
            urgency_score = 70

        task = {
            "action_type": "Message",
            "title": next_step,
            "due_at": due_dt.isoformat(),
            "stage": "Active Outreach",
            "last_activity_at": now.isoformat(),
            "urgency_score": urgency_score,
            "attempts": 0,
            "intent_score": intent_score,
            "stakeholders_detected": stakeholder_roles,
        }

        save_task(lead_id, task)

        increment_metric("messages")

        print(f"\nAuto discovery follow-up task created for lead '{lead_id}'.")
        print(f"Next step: {task['title']}")
        print(f"Due: {format_due_human(task['due_at'])}")
        print(f"Close Strength: {close_strength}")
        print(f"Intent Score: {intent_score}")
        if calendar_hold_recommended:
            print("⚠ Calendar hold recommended to protect momentum.")

        coaching_entry = {
            "lead_id": lead_id,
            "timestamp": now.isoformat(),
            "control_score": intent_score,
            "momentum": parsed.get("momentum"),
            "close_strength": close_strength,
            "stakeholder_gaps": stakeholder_roles,
            "top_adjustment": close_feedback
        }

        save_coaching_log(coaching_entry)

        print(f"Coaching Insight: {close_feedback}")
        print("Coaching log saved.\n")
        return

    print("No recognized action flag. Run with --morning, --brief, --today, --eod, --new-prospect, --post-meeting, --discovery, --auto-discovery, --attempt, --next-move, or --make-podcast.")


if __name__ == "__main__":
    main()